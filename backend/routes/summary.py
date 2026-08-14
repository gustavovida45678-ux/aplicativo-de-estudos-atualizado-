from fastapi import APIRouter, UploadFile, File, Form, HTTPException
from typing import Optional, List
from pydantic import BaseModel, Field
import asyncio
import logging
import json
import re
import uuid
from datetime import datetime, timezone

logger = logging.getLogger(__name__)
router = APIRouter()

# Groq free tier do llama-3.3-70b-versatile tem limite de 12.000 tokens/min (TPM).
# 8k caracteres + prompt JSON (~800 tokens) ficam perto de ~7.5k tokens, com margem segura.
MAX_MATERIAL_CHARS = 8000


async def _chat_with_rate_limit_retry(chat_service, message_fn, **kwargs):
    """Chama a IA com nova tentativa quando o provedor responde rate limit (ex: TPM do Groq).

    message_fn(max_chars) gera o prompt com o texto truncado no tamanho dado.
    A cada rate limit, o texto e reduzido em 50% antes de tentar de novo
    (8000 -> 4000 -> 2000 caracteres).
    """
    max_attempts = 3
    max_chars = MAX_MATERIAL_CHARS
    for attempt in range(max_attempts):
        try:
            message = message_fn(max_chars)
            return await chat_service.chat(message=message, **kwargs)
        except Exception as e:
            msg = str(e).lower()
            is_rate_limit = "rate limit" in msg or "rate_limit" in msg or "429" in msg
            if not is_rate_limit or attempt == max_attempts - 1:
                raise
            max_chars = max_chars * 50 // 100
            wait = 15 * (attempt + 1)
            logger.warning(
                f"Rate limit do provedor de IA (tentativa {attempt + 1}/{max_attempts}); "
                f"reduzindo material para {max_chars} chars, aguardando {wait}s..."
            )
            await asyncio.sleep(wait)


async def _chat_with_emergent_fallback(chat_service, message_fn, **kwargs):
    """Tenta o provider principal; se falhar (rate limit, timeout, TPM), tenta o EMERGENT."""
    try:
        return await _chat_with_rate_limit_retry(chat_service, message_fn, **kwargs)
    except Exception as e:
        from backend.services.providers import ProviderType
        logger.warning(
            f"Provider {kwargs.get('provider_type')} falhou ({type(e).__name__}: {str(e)[:200]}); "
            f"tentando EMERGENT..."
        )
        try:
            kwargs = dict(kwargs)
            kwargs["provider_type"] = ProviderType.EMERGENT
            return await _chat_with_rate_limit_retry(chat_service, message_fn, **kwargs)
        except Exception as e2:
            logger.warning(f"EMERGENT tambem falhou ({type(e2).__name__}: {str(e2)[:200]})")
            raise e2


# ── Models ──────────────────────────────────────────────────────────────────

class SummaryRequest(BaseModel):
    text: str = ""
    mode: str = "full"  # full | quick | study_focus
    exercise_count: int = 10
    difficulty: str = "misto"  # facil | medio | dificil | misto
    is_simulado: bool = False


class ExerciseAnswer(BaseModel):
    exercise_index: int
    answer: str


class SubmitAnswersRequest(BaseModel):
    summary_id: str
    answers: List[ExerciseAnswer]


# ── Helpers ─────────────────────────────────────────────────────────────────

def _extract_json(text: str):
    """Extract the first valid JSON object from text, ignoring markdown fences."""
    text = (text or "").strip()
    text = re.sub(r"^```(?:json)?\s*", "", text)
    text = re.sub(r"\s*```$", "", text)
    start = text.find("{")
    if start == -1:
        return None
    depth = 0
    in_str = False
    esc = False
    for i in range(start, len(text)):
        c = text[i]
        if in_str:
            if esc:
                esc = False
            elif c == "\\":
                esc = True
            elif c == '"':
                in_str = False
            continue
        if c == '"':
            in_str = True
        elif c == "{":
            depth += 1
        elif c == "}":
            depth -= 1
            if depth == 0:
                return text[start : i + 1]
    return None


def _extract_text_from_bytes(content: bytes, filename: str) -> str:
    """Extract text from file bytes (PDF, DOCX, TXT, images via OCR)."""
    lower = filename.lower()

    # TXT / Markdown
    if lower.endswith((".txt", ".md", ".csv", ".py", ".js", ".java", ".c", ".cpp", ".h")):
        try:
            return content.decode("utf-8", errors="replace")
        except Exception:
            return content.decode("latin-1", errors="replace")

    # PDF
    if lower.endswith(".pdf"):
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(stream=content, filetype="pdf")
            parts = []
            for page in doc:
                parts.append(page.get_text())
            text = "\n".join(parts).strip()
            if len(text) > 50:
                return text
        except Exception as e:
            logger.warning(f"PyMuPDF failed: {e}")

        # fallback: pdfplumber
        try:
            import pdfplumber
            import io
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                parts = [p.extract_text() or "" for p in pdf.pages]
            text = "\n".join(parts).strip()
            if len(text) > 50:
                return text
        except Exception as e:
            logger.warning(f"pdfplumber failed: {e}")

        # OCR fallback for scanned PDFs
        try:
            import fitz
            from PIL import Image as PILImage
            import pytesseract
            import io as _io
            doc = fitz.open(stream=content, filetype="pdf")
            parts = []
            for page in doc:
                pix = page.get_pixmap(dpi=200)
                img = PILImage.open(_io.BytesIO(pix.tobytes("png")))
                text = pytesseract.image_to_string(img, lang="eng+por")
                if text.strip():
                    parts.append(text)
            return "\n".join(parts).strip()
        except Exception as e:
            logger.warning(f"OCR PDF failed: {e}")
        return ""

    # DOCX
    if lower.endswith(".docx"):
        try:
            from docx import Document
            import io
            doc = Document(io.BytesIO(content))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception as e:
            logger.warning(f"docx failed: {e}")
        return ""

    # PPTX
    if lower.endswith(".pptx"):
        try:
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(content))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            txt = para.text.strip()
                            if txt:
                                parts.append(txt)
            return "\n".join(parts)
        except Exception as e:
            logger.warning(f"pptx failed: {e}")
        return ""

    # Images (OCR)
    if lower.endswith((".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff")):
        try:
            from PIL import Image as PILImage
            import pytesseract
            import io
            img = PILImage.open(io.BytesIO(content))
            return pytesseract.image_to_string(img, lang="eng+por", config="--psm 6").strip()
        except Exception as e:
            logger.warning(f"OCR image failed: {e}")
        return ""

    # DOC (legacy)
    if lower.endswith(".doc"):
        try:
            import subprocess
            import tempfile
            with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as tmp:
                tmp.write(content)
                tmp_path = tmp.name
            result = subprocess.run(
                ["textutil", "-convert", "txt", tmp_path, "-output", tmp_path + ".txt"],
                capture_output=True, timeout=30
            )
            with open(tmp_path + ".txt", "r", errors="replace") as f:
                return f.read()
        except Exception as e:
            logger.warning(f"doc conversion failed: {e}")
        return ""

    return ""


# ── Endpoints ───────────────────────────────────────────────────────────────

@router.post("/analyze")
async def analyze_material(
    files: List[UploadFile] = File(...),
    mode: str = Form("full"),
    exercise_count: int = Form(10),
    difficulty: str = Form("misto"),
    is_simulado: bool = Form(False),
):
    """Upload files, extract text, and generate a full summary + exercises via AI."""
    try:
        logger.info(f"📚 Summary analysis request: {len(files)} file(s), mode={mode}")

        # ── 1. Extract text from all files ──────────────────────────────────
        all_text_parts = []
        file_names = []
        for f in files:
            content = await f.read()
            text = _extract_text_from_bytes(content, f.filename or "unknown")
            if text:
                all_text_parts.append(text)
                file_names.append(f.filename or "unknown")
                logger.info(f"  ✅ {f.filename}: {len(text)} chars extracted")
            else:
                logger.warning(f"  ⚠️ {f.filename}: no text extracted")

        combined_text = "\n\n".join(all_text_parts)

        if not combined_text.strip():
            raise HTTPException(
                status_code=400,
                detail="Não foi possível extrair texto dos arquivos enviados. Verifique os formatos."
            )

        # Truncar para caber no TPM do Groq free tier (~12k tokens/min)
        if len(combined_text) > MAX_MATERIAL_CHARS:
            combined_text = combined_text[:MAX_MATERIAL_CHARS] + "\n\n[... texto truncado para processamento ...]"

        # ── 3. Call AI for summary ──────────────────────────────────────────
        from backend.services.chat_service import chat_service
        from backend.services.providers import ProviderType

        logger.info("🤖 Generating summary via AI...")
        summary_prompt_fn = lambda mc: (
            _build_quick_summary_prompt(combined_text[:mc])
            if mode == "quick"
            else _build_study_focus_prompt(combined_text[:mc])
            if mode == "study_focus"
            else _build_full_summary_prompt(combined_text[:mc])
        )
        summary_response = await _chat_with_emergent_fallback(
            chat_service,
            summary_prompt_fn,
            provider_type=ProviderType.GROQ,
            system_prompt="Você é um professor expert que analisa materiais de estudo e cria resumos inteligentes, completos e didáticos. Retorne APENAS JSON válido, sem markdown.",
            temperature=0.3,
            max_tokens=8192,
        )
        summary_text = summary_response.get("content", "")
        logger.info(f"✅ Summary AI response: {len(summary_text)} chars")

        summary_json = _extract_json(summary_text)
        if summary_json:
            summary_data = json.loads(summary_json)
        else:
            summary_data = _build_fallback_summary(combined_text, file_names)

        # ── 4. Call AI for exercises ────────────────────────────────────────
        logger.info("🤖 Generating exercises via AI...")
        exercises_response = await _chat_with_emergent_fallback(
            chat_service,
            lambda mc: _build_exercises_prompt(
                combined_text[:mc], exercise_count, difficulty, is_simulado
            ),
            provider_type=ProviderType.GROQ,
            system_prompt="Você é um professor expert em criar exercícios educacionais de altíssima qualidade. Retorne APENAS JSON válido, sem markdown.",
            temperature=0.4,
            max_tokens=8192,
        )
        exercises_text = exercises_response.get("content", "")
        logger.info(f"✅ Exercises AI response: {len(exercises_text)} chars")

        exercises_json = _extract_json(exercises_text)
        if exercises_json:
            exercises_data = json.loads(exercises_json)
        else:
            exercises_data = {"exercises": []}

        # ── 5. Build final response ─────────────────────────────────────────
        summary_id = str(uuid.uuid4())
        result = {
            "summary_id": summary_id,
            "files": file_names,
            "total_chars": len(combined_text),
            "summary": summary_data,
            "exercises": exercises_data.get("exercises", []),
            "mode": mode,
            "difficulty": difficulty,
            "is_simulado": is_simulado,
            "created_at": datetime.now(timezone.utc).isoformat(),
        }

        logger.info(f"✅ Summary complete: {len(result['summary'].get('topics', []))} topics, {len(result['exercises'])} exercises")
        return result

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"❌ Summary error: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Erro ao analisar material: {str(e)}")


@router.post("/exercises")
async def generate_more_exercises(
    content_text: str = Form(""),
    exercise_count: int = Form(10),
    difficulty: str = Form("misto"),
    is_simulado: bool = Form(False),
):
    """Generate additional exercises from previously analyzed content."""
    try:
        logger.info(f"📝 Extra exercises: count={exercise_count}, difficulty={difficulty}")

        if len(content_text) > MAX_MATERIAL_CHARS:
            content_text = content_text[:MAX_MATERIAL_CHARS]

        from backend.services.chat_service import chat_service
        from backend.services.providers import ProviderType

        response = await _chat_with_emergent_fallback(
            chat_service,
            lambda mc: _build_exercises_prompt(
                content_text[:mc], exercise_count, difficulty, is_simulado
            ),
            provider_type=ProviderType.GROQ,
            system_prompt="Você é um professor expert em criar exercícios educacionais. Retorne APENAS JSON válido, sem markdown.",
            temperature=0.4,
            max_tokens=8192,
        )
        response_text = response.get("content", "")

        json_str = _extract_json(response_text)
        if json_str:
            data = json.loads(json_str)
            return {"exercises": data.get("exercises", [])}

        return {"exercises": []}

    except Exception as e:
        logger.error(f"❌ Extra exercises error: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/validate-exercise")
async def validate_exercise(
    question: str = Form(""),
    correct_answer: str = Form(""),
    topic: str = Form(""),
    material_text: str = Form(""),
):
    """Validate that an exercise is grounded in the material."""
    try:
        from backend.services.chat_service import chat_service
        from backend.services.providers import ProviderType

        prompt = f"""Valide se este exercício está correto e fundamentado no material.

EXERCÍCIO: {question}
RESPOSTA CORRETA: {correct_answer}
TÓPICO: {topic}

MATERIAL (trecho):
{material_text[:3000]}

Responda APENAS com JSON:
{{"valid": true/false, "reason": "motivo"}}"""

        response = await chat_service.chat(
            message=prompt,
            provider_type=ProviderType.GROQ,
            system_prompt="Você é um validador de exercícios educacionais. Retorne APENAS JSON válido.",
            temperature=0.1,
            max_tokens=512,
        )
        text = response.get("content", "")
        json_str = _extract_json(text)
        if json_str:
            return json.loads(json_str)
        return {"valid": True, "reason": "Could not validate"}

    except Exception as e:
        logger.warning(f"Validation failed: {e}")
        return {"valid": True, "reason": "Validation service unavailable"}


# ── Prompt Builders ─────────────────────────────────────────────────────────

def _build_full_summary_prompt(text: str) -> str:
    return f"""Analise TODO o material de estudo abaixo e crie um RESUMO INTELIGENTE COMPLETO.

MATERIAL:
{text}

Responda EXATAMENTE com este JSON (sem markdown, sem ```):

{{
  "title": "Título do material",
  "general_summary": "Resumo geral do conteúdo em 3-5 parágrafos explicando o que o material aborda de forma clara e didática",
  "discipline": "Disciplina/área identificada",
  "topics": [
    {{
      "name": "Nome do Tópico",
      "priority": "alta",
      "explanation": "Explicação didática do tópico como se fosse um professor ensinando",
      "key_concepts": ["Conceito 1", "Conceito 2"],
      "formulas": ["Fórmula 1 = ..."],
      "examples": ["Exemplo prático 1"],
      "common_errors": ["Erro comum 1"],
      "observations": ["Observação importante"],
      "subtopics": [
        {{
          "name": "Subtítulo",
          "explanation": "Explicação do subtópico"
        }}
      ]
    }}
  ],
  "keywords": [
    {{
      "word": "Palavra-chave",
      "definition": "Definição clara",
      "explanation": "Explicação detalhada",
      "related_topics": ["Tópico relacionado"],
      "example": "Exemplo de uso"
    }}
  ],
  "concept_map": {{
    "root": "Disciplina",
    "branches": [
      {{
        "name": "Tópico Principal",
        "children": ["Subtópico 1", "Subtópico 2"]
      }}
    ]
  }},
  "formulas_summary": ["Fórmula 1: ...", "Fórmula 2: ..."],
  "study_tips": ["Dica 1 para estudar", "Dica 2"],
  "estimated_study_time": "X horas",
  "priority_analysis": [
    {{
      "topic": "Nome do tópico",
      "priority": "alta",
      "reason": "Por que é importante"
    }}
  ]
}}

REGRAS:
1. Extraia APENAS informações presentes no material
2. Priorize tópicos: alta (fundamental), média (relevante), baixa (complementar)
3. Cada palavra-chave deve ter definição e relação com o conteúdo
4. Seja didático como um professor particular
5. Retorne APENAS o JSON válido"""


def _build_quick_summary_prompt(text: str) -> str:
    return f"""Crie um RESUMO RÁPIDO do material abaixo.

MATERIAL:
{text}

Responda EXATAMENTE com este JSON (sem markdown, sem ```):

{{
  "title": "Título do material",
  "general_summary": "Resumo em 2 parágrafos",
  "discipline": "Disciplina",
  "topics": [
    {{
      "name": "Tópico",
      "priority": "alta",
      "explanation": "Explicação curta",
      "key_concepts": ["Conceito"],
      "formulas": [],
      "examples": [],
      "common_errors": [],
      "observations": [],
      "subtopics": []
    }}
  ],
  "keywords": [
    {{"word": "Palavra", "definition": "Definição", "explanation": "Explicação", "related_topics": [], "example": ""}}
  ],
  "concept_map": {{"root": "Disciplina", "branches": []}},
  "formulas_summary": [],
  "study_tips": [],
  "estimated_study_time": "X horas",
  "priority_analysis": []
}}

Retorne APENAS o JSON válido."""


def _build_study_focus_prompt(text: str) -> str:
    return f"""Analise o material e responda: "Se eu tivesse apenas este material para estudar para uma prova, o que eu deveria aprender primeiro?"

MATERIAL:
{text}

Responda EXATAMENTE com este JSON (sem markdown, sem ```):

{{
  "title": "Título do material",
  "general_summary": "Visão geral em 2 parágrafos focada no que mais cai em prova",
  "discipline": "Disciplina",
  "topics": [
    {{
      "name": "Tópico",
      "priority": "alta",
      "explanation": "Por que este tópico é o mais importante para prova",
      "key_concepts": ["Conceito essencial"],
      "formulas": ["Fórmula que mais cai"],
      "examples": ["Exemplo típico de prova"],
      "common_errors": ["Erro que leva zero"],
      "observations": ["Cai todo ano"],
      "subtopics": []
    }}
  ],
  "keywords": [
    {{"word": "Palavra", "definition": "Definição", "explanation": "Explicação", "related_topics": [], "example": ""}}
  ],
  "concept_map": {{"root": "Disciplina", "branches": []}},
  "formulas_summary": [],
  "study_tips": ["Dica de prova 1", "Dica de prova 2"],
  "estimated_study_time": "X horas",
  "priority_analysis": []
}}

Retorne APENAS o JSON válido.
REGRAS DE TIPO: TODO campo de texto (title, general_summary, name, explanation, word, definition, example, root, study_tips etc.) deve ser uma STRING de texto puro, nunca objeto, lista ou JSON aninhado. Campos de lista (key_concepts, formulas, examples, common_errors, observations, study_tips, priority_analysis) contêm APENAS strings. Se precisar enumerar algo dentro de um texto, use quebras de linha. NUNCA escreva literalmente "object Object" na resposta."""


def _build_exercises_prompt(text: str, count: int, difficulty: str, is_simulado: bool) -> str:
    diff_instruction = ""
    if difficulty == "facil":
        diff_instruction = "Todos os exercícios devem ser FÁCEIS (nível básico)."
    elif difficulty == "medio":
        diff_instruction = "Todos os exercícios devem ser de NÍVEL MÉDIO."
    elif difficulty == "dificil":
        diff_instruction = "Todos os exercícios devem ser DIFÍCEIS (avançados)."
    else:
        diff_instruction = "Distribua os exercícios: 30% fáceis, 40% médios, 30% difíceis."

    mode_instruction = ""
    if is_simulado:
        mode_instruction = "Crie uma PROVA/SIMULADO completo com questões progressivas (início fácil, meio intermediário, fim difícil). Inclua questões de múltipla escolha, verdadeiro/falso, complete a frase, resposta curta e discursiva."
    else:
        mode_instruction = """Varie os tipos de questões:
- 40% múltipla escolha (A, B, C, D)
- 15% verdadeiro ou falso
- 15% complete a frase
- 10% resposta curta
- 10% associação de conceitos
- 10% questões discursivas/problemas práticos"""

    return f"""Crie {count} EXERCÍCIOS baseados EXCLUSIVAMENTE no material abaixo.

MATERIAL:
{text}

{diff_instruction}

{mode_instruction}

DISTRIBUIÇÃO POR IMPORTÂNCIA:
- 50% dos exercícios → tópicos de ALTA prioridade
- 30% → tópicos de MÉDIA prioridade
- 20% → tópicos complementares

Para cada exercício, inclua:
- Pergunta clara e objetiva
- Alternativas (para múltipla escolha) ou resposta esperada
- Resposta correta
- Explicação DIDÁTICA ensinando POR QUE a resposta está correta
- Passo a passo da resolução
- Tópico relacionado

Responda EXATAMENTE com este JSON (sem markdown, sem ```):

{{
  "exercises": [
    {{
      "question": "Enunciado completo do exercício",
      "type": "multipla_escolha",
      "options": ["A) ...", "B) ...", "C) ...", "D) ..."],
      "correct_answer": "A",
      "explanation": "Explicação didática de POR QUE esta resposta está correta, ensinando o conceito",
      "solution_steps": [
        "Passo 1: ...",
        "Passo 2: ...",
        "Passo 3: ..."
      ],
      "topic": "Tópico relacionado",
      "difficulty": "Médio",
      "concept_used": "Conceito principal testado"
    }}
  ]
}}

REGRAS:
1. Exercícios devem estar RELACIONADOS ao material
2. Não invente conteúdo ausente no material
3. VARIE as respostas corretas (não use sempre a mesma letra)
4. Explique o PORQUÊ, não apenas o QUÊ
5. Retorne APENAS o JSON válido
6. IMPORTANTE: TODO campo de texto (question, correct_answer, explanation, topic, difficulty, concept_used, cada item de options e solution_steps) deve ser uma STRING de texto puro, nunca objeto, lista ou JSON aninhado. options e solution_steps contêm APENAS strings. NUNCA escreva literalmente "object Object" na resposta."""


def _build_fallback_summary(text: str, file_names: list) -> dict:
    """Build a basic summary structure when AI fails."""
    return {
        "title": "Resumo do Material",
        "general_summary": text[:1000] + ("..." if len(text) > 1000 else ""),
        "discipline": "Geral",
        "topics": [
            {
                "name": "Conteúdo Principal",
                "priority": "alta",
                "explanation": text[:500],
                "key_concepts": [],
                "formulas": [],
                "examples": [],
                "common_errors": [],
                "observations": [],
                "subtopics": [],
            }
        ],
        "keywords": [],
        "concept_map": {"root": "Material", "branches": []},
        "formulas_summary": [],
        "study_tips": [],
        "estimated_study_time": "1-2 horas",
        "priority_analysis": [],
    }
